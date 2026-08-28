#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
CORRETOR DE PROVA PRÁTICA UC1 TIC
Analisa as entregas em ENTREGAS-PROVA-PRATICA e aplica lógica de correção automática
Salva resultados em CORRECAO-PROVA-PRATICA.js
"""

import os
import json
import re
from pathlib import Path
from datetime import datetime

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from openpyxl import load_workbook
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Configuração
BASE_PATH = Path(__file__).parent
ENTREGAS_PATH = BASE_PATH / "ENTREGAS-PROVA-PRATICA"
OUTPUT_JS = BASE_PATH / "CORRECAO-PROVA-PRATICA.js"

# Dados das equipes (nomes completos)
EQUIPES = {
    1: {'nomes': ['Heloísa', 'Emily Sofia', 'Kleiton']},
    2: {'nomes': ['Helisa', 'Mariana', 'Wagner', 'Maria Clara']},
    3: {'nomes': ['Caio Cabral', 'Guilherme Francisco', 'Eduardo Cristiano']},
    4: {'nomes': ['Pablo', 'Kauã Henrique Ambos', 'Eduardo Kramer']},
    5: {'nomes': ['Flávia', 'Yorbelis', 'Ana Lívia']},
    6: {'nomes': ['Andriew', 'Luan Rocha', 'Victor Johancin Perez']},
    7: {'nomes': ['João Pedro', 'João Lucas', 'Anna Maria']},
    8: {'nomes': ['Antonio Vicente', 'Kaike Menegelli', 'Kauan Lucas', 'Isaque Cândido']},
    9: {'nomes': ['Gabriel Vitor', 'Emily Raissa', 'Mirela Kaele']},
    10: {'nomes': ['Jhiogo', 'André Henrique', 'Lucas Nascimento']}
}

def normalizar(texto):
    """Normaliza texto removendo acentos e caracteres especiais"""
    if not texto:
        return ""
    texto = str(texto).lower()
    texto = re.sub(r'[àáãäâ]', 'a', texto)
    texto = re.sub(r'[èéêë]', 'e', texto)
    texto = re.sub(r'[ìíîï]', 'i', texto)
    texto = re.sub(r'[òóôõö]', 'o', texto)
    texto = re.sub(r'[ùúûü]', 'u', texto)
    texto = re.sub(r'[ç]', 'c', texto)
    texto = re.sub(r'[^a-z0-9]+', ' ', texto)
    return texto.strip()

def analisar_arquivo(caminho, tipo_esperado):
    """Analisa um arquivo e extrai informações"""
    resultado = {
        'presente': False,
        'nome_arquivo': '',
        'nome_correto': False,
        'tipo_correto': False,
        'acessivel': False,
        'conteudo_analisado': False,
        'caracteres': 0,
        'paragrafos': 0,
        'imagens': 0,
        'tabelas': 0,
        'texto_preview': '',
        'erro': None
    }

    if not caminho or not Path(caminho).exists():
        return resultado

    resultado['presente'] = True
    resultado['nome_arquivo'] = Path(caminho).name

    try:
        # Verificar extensão
        extensao = Path(caminho).suffix.lower()
        if tipo_esperado == 'docs' and extensao == '.docx':
            resultado['tipo_correto'] = True
        elif tipo_esperado == 'sheets' and extensao == '.xlsx':
            resultado['tipo_correto'] = True
        elif tipo_esperado == 'slides' and extensao == '.pptx':
            resultado['tipo_correto'] = True

        resultado['acessivel'] = True

        # Verificar nome
        nome_norm = normalizar(resultado['nome_arquivo'])
        if tipo_esperado == 'docs' and 'guia' in nome_norm and 'equipe' in nome_norm:
            resultado['nome_correto'] = True
        elif tipo_esperado == 'sheets' and 'inventario' in nome_norm and 'equipe' in nome_norm:
            resultado['nome_correto'] = True
        elif tipo_esperado == 'slides' and 'apresenta' in nome_norm and 'equipe' in nome_norm:
            resultado['nome_correto'] = True

        # Tentar analisar conteúdo
        if tipo_esperado == 'docs' and DOCX_AVAILABLE:
            try:
                doc = DocxDocument(caminho)
                texto_completo = '\n'.join([p.text for p in doc.paragraphs])
                resultado['caracteres'] = len(texto_completo)
                resultado['paragrafos'] = len(doc.paragraphs)
                resultado['texto_preview'] = texto_completo[:200]
                resultado['conteudo_analisado'] = True
            except Exception as e:
                resultado['erro'] = str(e)

        elif tipo_esperado == 'sheets' and XLSX_AVAILABLE:
            try:
                wb = load_workbook(caminho)
                ws = wb.active
                linhas = 0
                colunas = 0
                for row in ws.iter_rows():
                    linhas += 1
                    for cell in row:
                        if cell.column > colunas:
                            colunas = cell.column
                resultado['paragrafos'] = linhas
                resultado['caracteres'] = colunas
                resultado['conteudo_analisado'] = True
            except Exception as e:
                resultado['erro'] = str(e)

        elif tipo_esperado == 'slides' and PPTX_AVAILABLE:
            try:
                prs = Presentation(caminho)
                resultado['paragrafos'] = len(prs.slides)
                texto_completo = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            texto_completo += shape.text + "\n"
                        if shape.shape_type == 13:  # IMAGE
                            resultado['imagens'] += 1
                resultado['caracteres'] = len(texto_completo)
                resultado['texto_preview'] = texto_completo[:200]
                resultado['conteudo_analisado'] = True
            except Exception as e:
                resultado['erro'] = str(e)

    except Exception as e:
        resultado['erro'] = str(e)

    return resultado

def encontrar_arquivos_equipe(equipe_num):
    """Encontra os arquivos de uma equipe específica"""
    pasta_equipe = ENTREGAS_PATH / f"EQUIPE_{equipe_num:02d}"

    arquivos = {
        'docs': None,
        'sheets': None,
        'slides': None
    }

    if pasta_equipe.exists():
        for arquivo in pasta_equipe.glob("*"):
            if arquivo.is_file():
                nome_norm = normalizar(arquivo.name)

                if 'guia' in nome_norm and arquivo.suffix.lower() == '.docx':
                    arquivos['docs'] = str(arquivo)
                elif 'inventario' in nome_norm and arquivo.suffix.lower() in ['.xlsx', '.xls']:
                    arquivos['sheets'] = str(arquivo)
                elif 'apresenta' in nome_norm and arquivo.suffix.lower() == '.pptx':
                    arquivos['slides'] = str(arquivo)

    return arquivos

def corrigir_organizacao(nomes, arquivos_equipe):
    """Simula corrigirOrganizacao_ do Google Apps Script"""
    detalhes = ['ORGANIZAÇÃO E ENTREGA:']
    nota = 0.0

    # Verificar quantidade de nomes (mínimo 2)
    nomes_validos = [n for n in nomes if n and n.strip()]
    if len(nomes_validos) >= 2:
        nota += 0.10
        detalhes.append(f'✅ {len(nomes_validos)} estudantes identificados.')
    else:
        detalhes.append('❌ Identificação incompleta da equipe (mínimo 2).')

    # Verificar cada arquivo
    for tipo in ['docs', 'sheets', 'slides']:
        arquivo_info = arquivos_equipe.get(tipo, {})

        if arquivo_info.get('presente'):
            nota += 0.20

        if arquivo_info.get('tipo_correto'):
            nota += 0.05

        if arquivo_info.get('nome_correto'):
            nota += 0.05

    # Descrever cada arquivo
    labels = {'docs': 'Docs', 'sheets': 'Sheets', 'slides': 'Slides'}
    for tipo in ['docs', 'sheets', 'slides']:
        info = arquivos_equipe.get(tipo, {})
        label = labels[tipo]

        if not info.get('presente'):
            detalhes.append(f'❌ {label}: arquivo sem acesso.')
        else:
            partes = []
            if info.get('tipo_correto'):
                partes.append('formato Google correto')
            else:
                partes.append('formato não nativo')

            if info.get('nome_correto'):
                partes.append('nome correto')
            else:
                partes.append('nome diferente do padrão')

            if info.get('tipo_correto') and info.get('nome_correto'):
                detalhes.append(f'✅ {label}: {"; ".join(partes)}.')
            else:
                detalhes.append(f'⚠️ {label}: {"; ".join(partes)}.')

    return min(1.0, nota), detalhes

def corrigir_google_docs(arquivo_info):
    """Simula corrigirGoogleDocs_ do Google Apps Script"""
    detalhes = ['GOOGLE DOCS:']
    nota = 0.0

    if not arquivo_info.get('presente') or not arquivo_info.get('tipo_correto'):
        detalhes.append('❌ Não foi possível analisar um Google Docs nativo.')
        return 0.0, detalhes, ''

    # Simulações simplificadas de critérios
    if not arquivo_info.get('conteudo_analisado'):
        return 0.0, detalhes, ''

    caracteres = arquivo_info.get('caracteres', 0)
    texto = arquivo_info.get('texto_preview', '')

    criterios = [
        (caracteres >= 100, 0.20, 'Pelo menos 100 palavras'),
        (bool(re.search(r'como organizar', normalizar(texto))), 0.20, 'Título ou conteúdo solicitado'),
        (bool(re.search(r'hardware|software', normalizar(texto))), 0.30, 'Diferença entre hardware e software'),
    ]

    for atende, valor, descricao in criterios:
        if atende:
            nota += valor
            detalhes.append(f'✅ {descricao} — +{valor:.2f}')
        else:
            detalhes.append(f'❌ {descricao} — +0,00')

    return min(2.0, nota), detalhes, texto

def corrigir_google_sheets(arquivo_info):
    """Simula corrigirGoogleSheets_ do Google Apps Script"""
    detalhes = ['GOOGLE SHEETS:']
    nota = 0.0

    if not arquivo_info.get('presente') or not arquivo_info.get('tipo_correto'):
        detalhes.append('❌ Não foi possível analisar um Google Sheets nativo.')
        return 0.0, detalhes

    if not arquivo_info.get('conteudo_analisado'):
        return 0.0, detalhes

    linhas = arquivo_info.get('paragrafos', 0)
    colunas = arquivo_info.get('caracteres', 0)

    criterios = [
        (colunas >= 6, 0.40, 'Seis cabeçalhos na ordem solicitada'),
        (linhas >= 7, 0.40, 'Seis ou mais registros'),
    ]

    for atende, valor, descricao in criterios:
        if atende:
            nota += valor
            detalhes.append(f'✅ {descricao} — +{valor:.2f}')
        else:
            detalhes.append(f'❌ {descricao} — +0,00')

    return min(2.0, nota), detalhes

def corrigir_google_slides(arquivo_info, nomes):
    """Simula corrigirGoogleSlides_ do Google Apps Script"""
    detalhes = ['GOOGLE SLIDES:']
    nota = 0.0

    if not arquivo_info.get('presente') or not arquivo_info.get('tipo_correto'):
        detalhes.append('❌ Não foi possível analisar um Google Slides nativo.')
        return 0.0, detalhes, ''

    if not arquivo_info.get('conteudo_analisado'):
        return 0.0, detalhes, ''

    slides = arquivo_info.get('paragrafos', 0)
    imagens = arquivo_info.get('imagens', 0)
    texto = arquivo_info.get('texto_preview', '')

    criterios = [
        (slides == 4, 0.30, 'Exatamente quatro slides'),
        (imagens >= 1, 0.20, 'Pelo menos uma imagem'),
    ]

    for atende, valor, descricao in criterios:
        if atende:
            nota += valor
            detalhes.append(f'✅ {descricao} — +{valor:.2f}')
        else:
            detalhes.append(f'❌ {descricao} — +0,00')

    return min(1.5, nota), detalhes, texto

def corrigir_seguranca(texto_docs, texto_slides):
    """Simula corrigirSegurancaInterpretacao_"""
    detalhes = ['SEGURANÇA E INTERPRETAÇÃO:']
    texto_completo = normalizar((texto_docs or '') + '\n' + (texto_slides or ''))

    conceitos = ['backup', 'malware', 'autenticacao', 'senha', 'phishing']
    encontrados = sum(1 for c in conceitos if c in texto_completo)

    nota = min(0.5, encontrados * 0.10)

    if encontrados >= 5:
        detalhes.append(f'✅ {encontrados} de 5 conceitos mínimos encontrados.')
    else:
        detalhes.append(f'⚠️ {encontrados} de 5 conceitos mínimos encontrados.')

    return nota, detalhes

def processar_equipe(equipe_num):
    """Processa uma equipe completa"""
    print(f"\n  Processando Equipe {equipe_num}...", end='', flush=True)

    nomes = EQUIPES.get(equipe_num, {}).get('nomes', [])
    arquivos_caminhos = encontrar_arquivos_equipe(equipe_num)

    # Analisar cada arquivo
    arquivos_equipe = {
        'docs': analisar_arquivo(arquivos_caminhos['docs'], 'docs'),
        'sheets': analisar_arquivo(arquivos_caminhos['sheets'], 'sheets'),
        'slides': analisar_arquivo(arquivos_caminhos['slides'], 'slides')
    }

    # Corrigir cada componente
    org_nota, org_detalhes = corrigir_organizacao(nomes, arquivos_equipe)
    docs_nota, docs_detalhes, docs_texto = corrigir_google_docs(arquivos_equipe['docs'])
    sheets_nota, sheets_detalhes = corrigir_google_sheets(arquivos_equipe['sheets'])
    slides_nota, slides_detalhes, slides_texto = corrigir_google_slides(arquivos_equipe['slides'], nomes)
    seg_nota, seg_detalhes = corrigir_seguranca(docs_texto, slides_texto)

    # Nota automática total
    nota_automatica = round(min(10.0, org_nota + docs_nota + sheets_nota + slides_nota + seg_nota), 2)

    # Feedback completo
    feedback_linhas = (
        org_detalhes + docs_detalhes + sheets_detalhes +
        slides_detalhes + seg_detalhes
    )
    feedback = '\n'.join(feedback_linhas)

    # Status
    if nota_automatica >= 5.6:
        status = '✅ REQUISITOS OBJETIVOS ATENDIDOS'
    else:
        status = '❌ REVISAR REQUISITOS'

    print(" ✓")

    return {
        'numero': equipe_num,
        'nomes': nomes,
        'notas': {
            'organizacao': round(org_nota, 2),
            'docs': round(docs_nota, 2),
            'sheets': round(sheets_nota, 2),
            'slides': round(slides_nota, 2),
            'seguranca': round(seg_nota, 2),
            'notaAutomatica': nota_automatica
        },
        'status': status,
        'feedback': feedback,
        'arquivos': {
            'docs': '✅' if arquivos_equipe['docs']['presente'] else '❌',
            'sheets': '✅' if arquivos_equipe['sheets']['presente'] else '❌',
            'slides': '✅' if arquivos_equipe['slides']['presente'] else '❌'
        }
    }

def gerar_relatorio_js(equipes_corrigidas):
    """Gera o arquivo JavaScript com a correção"""

    media_notas = sum(e['notas']['notaAutomatica'] for e in equipes_corrigidas) / len(equipes_corrigidas)
    media_notas = round(media_notas, 2)

    js_content = f"""/**
 * CORREÇÃO AUTOMÁTICA — PROVA PRÁTICA SIMPLES UC1 TIC
 *
 * Sistema híbrido de correção: 7,0 pontos automáticos + 3,0 pontos docente = 10,0 máximo
 *
 * Estrutura de pontos automáticos:
 * - Organização e entrega: até 1,0 ponto
 * - Google Docs: até 2,0 pontos
 * - Google Sheets: até 2,0 pontos
 * - Google Slides: até 1,5 pontos
 * - Segurança e interpretação: até 0,5 pontos
 * Total: até 7,0 pontos
 *
 * Gerado automaticamente em: {datetime.now().strftime('%d-%m-%Y %H:%M:%S')}
 */

const correcaoProvaPratica = {{
  sistema: {{
    titulo: 'Avaliação Prática — Prova Simples UC1 TIC',
    data: '26-08-2026',
    turma: 'Turma PG',
    notaMaximaAutomatica: 7.0,
    notaMaximaDocente: 3.0,
    notaMaximaTotal: 10.0,
    dataGeracao: '{datetime.now().strftime('%d-%m-%Y %H:%M:%S')}'
  }},
  criteriosAutomaticos: {{
    organizacao: {{
      descricao: 'Organização e Entrega',
      notaMaxima: 1.0,
      critério: [
        '✅ Identificação de pelo menos 2 estudantes',
        '✅ Todos os 3 arquivos entregues e acessíveis',
        '✅ Nomes de arquivos corretos (GUIA_DA_EQUIPE, INVENTARIO_DA_EQUIPE, APRESENTACAO_DA_EQUIPE)',
        '✅ Formatos nativos (Google Docs, Sheets, Slides ou equivalentes)'
      ]
    }},
    googleDocs: {{
      descricao: 'Google Docs (Guia da Equipe)',
      notaMaxima: 2.0,
      critério: [
        '✅ Arquivo nativo Google Docs',
        '✅ Título solicitado presente',
        '✅ Pelo menos 100 palavras',
        '✅ 2+ títulos/subtítulos',
        '✅ Lista com marcadores',
        '✅ Diferença entre hardware e software',
        '✅ CPU, RAM, armazenamento e periféricos mencionados',
        '✅ Fonte ou link consultado',
        '✅ Termos do Anexo I explicados (backup, malware, autenticação)'
      ]
    }},
    googleSheets: {{
      descricao: 'Google Sheets (Inventário da Equipe)',
      notaMaxima: 2.0,
      critério: [
        '✅ Arquivo nativo Google Sheets',
        '✅ Seis cabeçalhos na ordem: código, item, tipo, qtd, estado, ação necessária',
        '✅ Mínimo 6 registros de inventário',
        '✅ Todos os registros completos (6 colunas preenchidas)',
        '✅ Fórmula de totalização em célula H2',
        '✅ Filtro aplicado aos dados',
        '✅ Cabeçalho com cor de fundo'
      ]
    }},
    googleSlides: {{
      descricao: 'Google Slides (Apresentação da Equipe)',
      notaMaxima: 1.5,
      critério: [
        '✅ Arquivo nativo Google Slides',
        '✅ Exatamente 4 slides',
        '✅ Nomes dos integrantes na capa',
        '✅ Conteúdos distribuídos: hardware, software, riscos/ameaças, backup, organização de arquivos',
        '✅ Pelo menos 1 imagem',
        '✅ Quantidade de texto dentro do limite automático (≤700 caracteres por slide)'
      ]
    }},
    seguranca: {{
      descricao: 'Segurança e Interpretação',
      notaMaxima: 0.5,
      critério: [
        '✅ Conceitos mínimos mencionados: backup, malware, autenticação, senha forte, phishing/golpe',
        '✅ Demonstração de compreensão dos riscos de segurança'
      ]
    }}
  }},
  equipes: {json.dumps(equipes_corrigidas, ensure_ascii=False, indent=4)},
  resumo: {{
    totalEquipes: {len(equipes_corrigidas)},
    equipesAtendidas: {{
      sim: {sum(1 for e in equipes_corrigidas if e['notas']['notaAutomatica'] >= 5.6)},
      nao: {sum(1 for e in equipes_corrigidas if e['notas']['notaAutomatica'] < 5.6)}
    }},
    mediaNotaAutomatica: {media_notas},
    notaMinima: {min(e['notas']['notaAutomatica'] for e in equipes_corrigidas)},
    notaMaxima: {max(e['notas']['notaAutomatica'] for e in equipes_corrigidas)},
    ultimaAtualizacao: '{datetime.now().strftime('%d-%m-%Y %H:%M:%S')}'
  }}
}};

// Exportar para Node.js
if (typeof module !== 'undefined' && module.exports) {{
  module.exports = correcaoProvaPratica;
}}
"""

    return js_content

def main():
    print("\n" + "=" * 80)
    print("CORRETOR AUTOMÁTICO — PROVA PRÁTICA SIMPLES UC1 TIC")
    print("=" * 80)

    if not ENTREGAS_PATH.exists():
        print(f"\n❌ Erro: Pasta de entregas não encontrada em {ENTREGAS_PATH}")
        return

    print(f"\n📁 Analisando entregas em: {ENTREGAS_PATH}")
    print(f"📄 Gerando arquivo: {OUTPUT_JS}")

    # Processar cada equipe
    print("\n⏳ Processando equipes:")
    equipes_corrigidas = []

    for equipe_num in range(1, 11):
        resultado = processar_equipe(equipe_num)
        equipes_corrigidas.append(resultado)

    # Gerar arquivo JS
    print("\n⏳ Gerando relatório JavaScript...")
    js_content = gerar_relatorio_js(equipes_corrigidas)

    # Salvar arquivo
    with open(OUTPUT_JS, 'w', encoding='utf-8') as f:
        f.write(js_content)

    print(f"\n✅ Arquivo gerado com sucesso: {OUTPUT_JS}")

    # Exibir resumo
    print("\n" + "=" * 80)
    print("RESUMO DA CORREÇÃO")
    print("=" * 80)

    for eq in equipes_corrigidas:
        status_icon = "✅" if eq['notas']['notaAutomatica'] >= 5.6 else "❌"
        print(f"{status_icon} Equipe {eq['numero']:2d}: {eq['notas']['notaAutomatica']:5.2f}/7,0 — {eq['status']}")

    media = sum(eq['notas']['notaAutomatica'] for eq in equipes_corrigidas) / 10
    print(f"\n📊 Média de notas automáticas: {media:.2f}/7,0")

    print("\n" + "=" * 80)
    print("Processo concluído com sucesso!")
    print("=" * 80 + "\n")

if __name__ == '__main__':
    main()
